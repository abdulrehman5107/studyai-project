import os
import io
import json
from flask import Flask, request, jsonify
from flask_cors import CORS
import pdfplumber
from docx import Document
from groq import Groq

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

GROQ_API_KEY = os.environ.get("GROQ_API_KEY")

client = Groq(api_key=GROQ_API_KEY)

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def chunk_text(text, max_chars=6000):
    if len(text) <= max_chars:
        return text
    return text[:max_chars]

def generate_study_materials(text):
    text = chunk_text(text)

    prompt = f"""You are an expert educator. Given the following lecture content, generate study materials.

Return ONLY a valid JSON object with NO extra text, no markdown, no code fences.

Lecture Content:
{text}

Return exactly this JSON structure:
{{
  "summary": ["point 1", "point 2", "point 3", "point 4", "point 5"],
  "notes": "# Topic Title\\n\\n## Section 1\\n\\nDetailed explanation here.\\n\\n- **Key term**: definition\\n\\n## Section 2\\n\\nMore explanation here.",
  "flashcards": [
    {{"question": "Q1", "answer": "A1"}},
    {{"question": "Q2", "answer": "A2"}},
    {{"question": "Q3", "answer": "A3"}},
    {{"question": "Q4", "answer": "A4"}},
    {{"question": "Q5", "answer": "A5"}},
    {{"question": "Q6", "answer": "A6"}},
    {{"question": "Q7", "answer": "A7"}},
    {{"question": "Q8", "answer": "A8"}}
  ],
  "mcqs": [
    {{
      "question": "Question text?",
      "options": ["A) option1", "B) option2", "C) option3", "D) option4"],
      "answer": "A) option1",
      "explanation": "Because..."
    }}
  ],
  "practice_questions": [
    {{"question": "Short Q1", "sample_answer": "Brief answer"}},
    {{"question": "Short Q2", "sample_answer": "Brief answer"}},
    {{"question": "Short Q3", "sample_answer": "Brief answer"}},
    {{"question": "Long Q4", "sample_answer": "A detailed multi-sentence answer explaining the concept thoroughly with examples."}},
    {{"question": "Long Q5", "sample_answer": "A detailed multi-sentence answer explaining the concept thoroughly with examples."}}
  ]
}}

Rules:
- Include exactly 10 MCQs.
- For notes: use \\n for newlines, ## for headings, **term** for bold. Write at least 400 words.
- Return only valid JSON, nothing else."""

    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": "You are an expert educator. You always respond with valid JSON only, no extra text."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.3,
        max_tokens=4000
    )

    raw = response.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()
    try:
        return json.loads(raw)
    except Exception:
        raw = " ".join(raw.splitlines())
        return json.loads(raw)

def handle_request():
    try:
        text = ""
        if "file" in request.files:
            file = request.files["file"]
            filename = file.filename.lower()
            if filename.endswith(".pdf"):
                text = extract_text_from_pdf(io.BytesIO(file.read()))
            elif filename.endswith(".docx"):
                text = extract_text_from_docx(io.BytesIO(file.read()))
            elif filename.endswith(".pptx"):
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file.read()))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text_frame.text + "\n"
            elif filename.endswith(".txt"):
                text = file.read().decode("utf-8")
            else:
                return jsonify({"error": "Unsupported file type. Use PDF, DOCX, or TXT."}), 400
        elif "text" in request.form:
            text = request.form["text"]
        elif request.is_json:
            data = request.get_json()
            text = data.get("text", "")
        else:
            return jsonify({"error": "No file or text provided."}), 400

        if not text.strip():
            return jsonify({"error": "Could not extract text from the file."}), 400

        result = generate_study_materials(text)
        return jsonify(result)

    except json.JSONDecodeError as e:
        return jsonify({"error": f"AI returned invalid JSON: {str(e)}"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate", methods=["POST", "OPTIONS"])
def generate():
    if request.method == "OPTIONS":
        return jsonify({}), 200
    return handle_request()

@app.route("/api/process", methods=["POST", "OPTIONS"])
def process():
    if request.method == "OPTIONS":
        return jsonify({}), 200
    return handle_request()

@app.route("/api/health", methods=["GET"])
@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
